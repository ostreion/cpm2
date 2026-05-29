#!/usr/bin/env python
"""Build the CPM2 / CPNext status-report & hand-off deck.

Re-runnable: reads the benchmark CSVs + final_report.md, regenerates the
matplotlib PNGs under assets/, then assembles CPM2_status_report.pptx.

Run from anywhere with the cpm2 env active:
    python notebooks/slides/build_deck.py

Every quantitative claim is anchored to a file under benchmarks/. No numbers
are invented. Slides where data is weak carry an explicit epistemic footnote.
No em dashes anywhere (user style rule): periods, commas, parentheses, hyphens.
"""
from __future__ import annotations

from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from matplotlib.patches import Patch

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

# --------------------------------------------------------------------------
# Paths
# --------------------------------------------------------------------------
HERE = Path(__file__).resolve().parent
REPO = HERE.parents[1]
ASSETS = HERE / "assets"
ASSETS.mkdir(parents=True, exist_ok=True)

NEGCTRL = REPO / "benchmarks" / "md_20260528_negcontrols"
OVERNIGHT = REPO / "benchmarks" / "md_20260526_overnight"
OUT_PPTX = HERE / "CPM2_status_report.pptx"

# --------------------------------------------------------------------------
# Theme
# --------------------------------------------------------------------------
INK = RGBColor(0x1A, 0x1A, 0x2E)        # near-black title ink
SLATE = RGBColor(0x3D, 0x40, 0x55)      # body text
MUTED = RGBColor(0x6B, 0x6F, 0x80)      # footnotes
ACCENT = RGBColor(0x2D, 0x6A, 0x8E)     # teal-blue accent
GOOD = RGBColor(0x2E, 0x7D, 0x52)       # green
WARN = RGBColor(0xB4, 0x4A, 0x1E)       # burnt-orange caution
BG = RGBColor(0xFB, 0xFB, 0xFD)
CARD = RGBColor(0xF1, 0xF2, 0xF6)

# matplotlib palette (consistent with deck)
MPL_COGNATE = "#2d6a8e"
MPL_NAT = "#2e7d52"
MPL_ALA = "#9aa0ad"
MPL_SCR = "#d98a3d"
MPL_WT = "#b44a1e"
MUTED_HEX = "#6b6f80"  # muted gray for in-figure annotations

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# --------------------------------------------------------------------------
# Data
# --------------------------------------------------------------------------
def load_data():
    ctrl = pd.read_csv(NEGCTRL / "results_running.csv")
    cog = pd.read_csv(OVERNIGHT / "results_running.csv")

    # Numbers come from the canonical backend (src/critique.py), the same
    # source the presentation notebook uses, so the deck and notebook cannot
    # silently diverge. critique reads the two CSVs above; we only map its
    # target keys to the deck's display labels here.
    import sys
    if str(REPO / "src") not in sys.path:
        sys.path.insert(0, str(REPO / "src"))
    import critique as C

    _key = {"mdm2_p53": "MDM2", "bclxl": "Bcl-xL",
            "14_3_3": "14-3-3", "cypA": "CypA"}
    _dvr = C.design_vs_reference().set_index("target")
    _spec = C.specificity_matrix()

    winners = {_key[k]: v for k, v in C.WINNERS.items()}
    design_dg = {_key[t]: round(float(r["design_dG"]), 1)
                 for t, r in _dvr.iterrows()}
    nat_dg = {_key[t]: round(float(r["reference_dG"]), 1)
              for t, r in _dvr.iterrows()}
    nat_label = {
        "MDM2": "p53 TAD",
        "Bcl-xL": "Bad (BH3)",
        "14-3-3": "phos-Cdc25C",
        "CypA": "HIV capsid",
    }
    nat_flag = {"Bcl-xL": "not converged"}  # honest caveat marker

    # ddG matrix (control - cognate), kcal/mol, from critique.specificity_matrix
    targets = ["MDM2", "Bcl-xL", "14-3-3", "CypA"]
    _piv = {}
    for _cond in ("alascan", "scrambled", "wrongtarget"):
        _col = {_key[row["target"]]: round(float(row["ddg_vs_cognate"]), 1)
                for _, row in _spec[_spec["condition"] == _cond].iterrows()}
        _piv[_cond] = [_col[t] for t in targets]
    ddg = pd.DataFrame(_piv, index=targets)

    return ctrl, cog, winners, design_dg, nat_dg, nat_label, nat_flag, ddg


# --------------------------------------------------------------------------
# Figure 1: design vs cognate-best vs natural-partner dG
# --------------------------------------------------------------------------
def make_bar_chart(design_dg, nat_dg, nat_label, nat_flag):
    targets = list(design_dg.keys())
    x = np.arange(len(targets))
    w = 0.38

    fig, ax = plt.subplots(figsize=(10.5, 5.4), dpi=150)
    des = [design_dg[t] for t in targets]
    nat = [nat_dg[t] for t in targets]

    # honest +-5 kcal/mol error bars (NOT the misleading MMPBSA 0.5-0.9)
    err = 5.0
    b1 = ax.bar(x - w / 2, des, w, label="Best design (cyclic)", color=MPL_COGNATE,
                yerr=err, capsize=5, error_kw=dict(ecolor="#555", lw=1.3))
    b2 = ax.bar(x + w / 2, nat, w, label="Natural partner (reference)", color=MPL_NAT,
                yerr=err, capsize=5, error_kw=dict(ecolor="#555", lw=1.3))

    ax.axhline(0, color="#888", lw=0.8)
    ax.set_ylabel(r"$\Delta G_{\mathrm{MMGBSA}}$ (kcal/mol, more negative = tighter)", fontsize=13)
    ax.set_xticks(x)
    ax.set_xticklabels(targets, fontsize=14)
    ax.tick_params(axis="y", labelsize=12)
    ax.set_title("Best designed cyclic peptide vs natural binding partner, per target",
                 fontsize=15, weight="bold", pad=14)
    ax.legend(fontsize=12, loc="lower right", frameon=False)

    # value labels
    for rects, vals, names in ((b1, des, [None] * 4), (b2, nat, [nat_label[t] for t in targets])):
        for r, v, nm in zip(rects, vals, names):
            ax.annotate(f"{v:.1f}", (r.get_x() + r.get_width() / 2, v),
                        textcoords="offset points", xytext=(0, -16 if v < 0 else 6),
                        ha="center", fontsize=11, color="white" if v < 0 else "#333",
                        weight="bold")
            if nm:
                ax.annotate(nm, (r.get_x() + r.get_width() / 2, 2),
                            ha="center", fontsize=9, color=MUTED_HEX, rotation=0)

    # flag Bcl-xL natural ref not converged (place to the right of the bar to avoid the value label)
    for i, t in enumerate(targets):
        if t in nat_flag:
            ax.annotate(f"* {nat_flag[t]}", (x[i] + w / 2 + 0.20, nat_dg[t] + 3),
                        ha="left", fontsize=9, color=MPL_WT, style="italic")

    ax.set_ylim(min(nat) - 16, 12)
    ax.spines[["top", "right"]].set_visible(False)
    fig.text(0.01, 0.01,
             "Error bars: honest +-5 kcal/mol (n=1 per cell, autocorrelated frames). "
             "Bcl-xL natural ref (*) not converged at 10 ns.",
             fontsize=9, color=MUTED_HEX)
    fig.tight_layout(rect=(0, 0.04, 1, 1))
    out = ASSETS / "fig_design_vs_natural.png"
    fig.savefig(out, dpi=150)
    plt.close(fig)
    return out


# --------------------------------------------------------------------------
# Figure 2: specificity ddG heatmap (centrepiece)
# --------------------------------------------------------------------------
def make_heatmap(ddg):
    rows = ddg.index.tolist()
    cols = ["alascan", "scrambled", "wrongtarget"]
    col_labels = ["Ala-scan\n(strip sidechains)", "Scrambled\n(permute sequence)",
                  "Wrong target\n(foreign pocket)"]
    M = ddg[cols].values

    fig, ax = plt.subplots(figsize=(8.8, 5.6), dpi=150)
    # green = strong discrimination (high ddG), white/red = no discrimination (low ddG)
    im = ax.imshow(M, cmap="YlGn", vmin=0, vmax=65, aspect="auto")

    ax.set_xticks(np.arange(len(cols)))
    ax.set_xticklabels(col_labels, fontsize=12)
    ax.set_yticks(np.arange(len(rows)))
    ax.set_yticklabels(rows, fontsize=13)
    ax.set_title(r"Negative-control specificity: $\Delta\Delta G$ vs cognate (kcal/mol)",
                 fontsize=15, weight="bold", pad=14)

    for i in range(len(rows)):
        for j in range(len(cols)):
            v = M[i, j]
            failed = v < 5  # within noise -> no discrimination
            txt = f"+{v:.1f}"
            color = MPL_WT if failed else ("#1a1a2e" if v < 35 else "white")
            ax.text(j, i, txt, ha="center", va="center", fontsize=14,
                    weight="bold", color=color)
            if failed:
                ax.text(j, i + 0.30, "no discrim.", ha="center", va="center",
                        fontsize=9, color=MPL_WT, style="italic")

    cb = fig.colorbar(im, ax=ax, fraction=0.046, pad=0.04)
    cb.set_label(r"$\Delta\Delta G$ (higher = better discrimination)", fontsize=10)
    fig.text(0.01, 0.015,
             "Green = control collapses (good: pipeline discriminates). Orange cell = control\n"
             "scores like cognate (failure of specificity). All cells +-5 kcal/mol.",
             fontsize=9, color=MUTED_HEX)
    fig.tight_layout(rect=(0.02, 0.08, 1, 1))
    out = ASSETS / "fig_specificity_heatmap.png"
    fig.savefig(out, dpi=150)
    plt.close(fig)
    return out




# --------------------------------------------------------------------------
# pptx helpers
# --------------------------------------------------------------------------
def _blank(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    # background fill
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def _textbox(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    return tb, tf


def _set(p, text, size, color=SLATE, bold=False, italic=False, align=PP_ALIGN.LEFT,
         font="Calibri"):
    p.alignment = align
    r = p.add_run()
    r.text = text
    f = r.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    f.name = font
    return r


def _accent_bar(slide):
    bar = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(0.12))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    return bar


def _title(slide, text, sub=None):
    _accent_bar(slide)
    tb, tf = _textbox(slide, Inches(0.55), Inches(0.30), Inches(12.2), Inches(1.0))
    _set(tf.paragraphs[0], text, 28, INK, bold=True)
    if sub:
        p = tf.add_paragraph()
        _set(p, sub, 14, ACCENT, italic=True)


def _footnote(slide, text, color=MUTED):
    tb, tf = _textbox(slide, Inches(0.55), Inches(6.95), Inches(12.2), Inches(0.45))
    _set(tf.paragraphs[0], "Epistemic flag: " + text, 11, color, italic=True)


def _bullets(slide, items, l, t, w, h, size=16, gap=6):
    tb, tf = _textbox(slide, l, t, w, h)
    first = True
    for it in items:
        if isinstance(it, tuple):
            txt, lvl, kw = it
        else:
            txt, lvl, kw = it, 0, {}
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = lvl
        p.space_after = Pt(gap)
        bullet = ("- " if lvl == 0 else "  . ")
        _set(p, bullet + txt, size - lvl * 1, kw.get("color", SLATE),
             bold=kw.get("bold", False), italic=kw.get("italic", False))
    return tb


def _card(slide, l, t, w, h, fill=CARD):
    box = slide.shapes.add_shape(1, l, t, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = RGBColor(0xDD, 0xDF, 0xE6)
    box.line.width = Pt(0.75)
    box.shadow.inherit = False
    return box


def _pic_fit(slide, path, l, t, max_w, max_h):
    from PIL import Image
    with Image.open(path) as im:
        iw, ih = im.size
    ar = iw / ih
    w = max_w
    h = Emu(int(w / ar))
    if h > max_h:
        h = max_h
        w = Emu(int(h * ar))
    # center horizontally within max_w region
    off = Emu(int((max_w - w) / 2))
    slide.shapes.add_picture(str(path), l + off, t, w, h)


# --------------------------------------------------------------------------
# Slides
# --------------------------------------------------------------------------
def build(prs, data, bar_png, heat_png):
    (ctrl, cog, winners, design_dg, nat_dg, nat_label, nat_flag, ddg) = data

    # ---- Slide 1: Title / hand-off ----
    s = _blank(prs)
    _accent_bar(s)
    tb, tf = _textbox(s, Inches(0.7), Inches(2.1), Inches(12), Inches(2.2))
    _set(tf.paragraphs[0], "CPM2 / CPNext: cyclic-peptide PPI-inhibitor pipeline", 36, INK, bold=True)
    p = tf.add_paragraph(); _set(p, "Status report and hand-off", 24, ACCENT, bold=True)
    p = tf.add_paragraph(); p.space_before = Pt(18)
    _set(p, "A candid evaluation of pipeline quality. What it can show, where it is weak, what is queued.", 16, SLATE, italic=True)

    tb, tf = _textbox(s, Inches(0.7), Inches(5.4), Inches(12), Inches(1.6))
    _set(tf.paragraphs[0], "Adrian Cipriani  ->  Niklas, Martin", 16, SLATE, bold=True)
    p = tf.add_paragraph(); _set(p, "Period: 2026-04-13 to ongoing.  Informal hand-off briefing.", 14, MUTED)
    p = tf.add_paragraph(); _set(p, "Scientific spine: negative-control MD batch (benchmarks/md_20260528_negcontrols).", 12, MUTED, italic=True)

    # ---- Slide 2: Question + pipeline diagram ----
    s = _blank(prs)
    _title(s, "The question, and the pipeline",
           "Can interface-mimicry seed competitive cyclic-peptide PPI inhibitors?")
    _bullets(s, [
        "Goal: design head-to-tail / disulfide cyclic peptides that block four protein-protein interfaces.",
        "Targets: MDM2/p53, Bcl-xL/Bad, 14-3-3/Cdc25C, CypA/HIV-capsid.",
        "Premise: borrow the backbone geometry of a known interface, then redesign the sequence.",
    ], Inches(0.55), Inches(1.6), Inches(12.2), Inches(1.7), size=16)

    # three-stage diagram
    stages = [
        ("cPEPmatch", "geometric backbone\nmatch to interface", ACCENT),
        ("Boltz-2", "fold-validate complex\n(ipTM, pLDDT, RMSD)", RGBColor(0x4A, 0x6B, 0x8A)),
        ("ProteinHunter", "sequence refinement\n(LigandMPNN + Boltz)", RGBColor(0x5E, 0x7C, 0x6E)),
    ]
    bx_w, bx_h = Inches(3.3), Inches(1.7)
    y = Inches(4.1)
    xs = [Inches(0.9), Inches(5.0), Inches(9.1)]
    for (name, desc, col), x in zip(stages, xs):
        box = _card(s, x, y, bx_w, bx_h, fill=col)
        tf = box.text_frame; tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        _set(tf.paragraphs[0], name, 20, RGBColor(0xFF, 0xFF, 0xFF), bold=True, align=PP_ALIGN.CENTER)
        p = tf.add_paragraph(); _set(p, desc, 13, RGBColor(0xEE, 0xEE, 0xF2), align=PP_ALIGN.CENTER)
    for x in xs[:-1]:
        ar = s.shapes.add_textbox(x + bx_w - Inches(0.1), y + Inches(0.5), Inches(0.9), Inches(0.7))
        _set(ar.text_frame.paragraphs[0], "->", 30, ACCENT, bold=True, align=PP_ALIGN.CENTER)
    _footnote(s, "Pipeline computes geometry + model confidence, not a thermodynamic score. "
                 "Affinity is assessed only later, by offline MD/MMGBSA on a shortlist.")

    # ---- Slide 3: Triage funnel ----
    s = _blank(prs)
    _title(s, "What the pipeline can show: the triage funnel",
           "MDM2/p53 v1 run as the worked example (data/runs/mdm2_p53_v1_...).")
    funnel = [
        ("37", "cPEPmatch backbone matches", ACCENT),
        ("37", "Boltz-2 folded poses (ipTM 0.83 to 0.96)", RGBColor(0x4A, 0x6B, 0x8A)),
        ("148", "ProteinHunter refined designs", RGBColor(0x5E, 0x7C, 0x6E)),
        ("3", "shortlist advanced to MD / MMGBSA", GOOD),
    ]
    y = Inches(1.75)
    widths = [Inches(11.5), Inches(9.5), Inches(7.5), Inches(4.0)]
    for (n, lab, col), w in zip(funnel, widths):
        l = Emu(int((SLIDE_W - w) / 2))
        box = _card(s, l, y, w, Inches(0.95), fill=col)
        tf = box.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        _set(tf.paragraphs[0], f"{n}   {lab}", 18, RGBColor(0xFF, 0xFF, 0xFF), bold=True, align=PP_ALIGN.CENTER)
        y = Emu(int(y) + int(Inches(1.18)))
    _bullets(s, [
        "Inspection layer: marimo triage view (ipTM vs peptide-drift plane, sized by pLDDT) plus a per-design drilldown.",
        "Enriched per-design table adds topology (head-to-tail / monocyclic / bicyclic), cPEPmatch fit-RMSD, peptide drift, and a topology audit that flags designs built with the wrong cyclization vs the DB.",
        "Shortlist picks the best cyclic design per topology family, so MD spends GPU only on a few coverage-spanning candidates.",
    ], Inches(0.55), Inches(6.0), Inches(12.2), Inches(1.0), size=13, gap=3)

    # ---- Slide 4: Design vs cognate/natural bar chart ----
    s = _blank(prs)
    _title(s, "Designs vs natural partners: MMGBSA per target",
           "Best designed cyclic peptide against the cognate natural binder.")
    _pic_fit(s, bar_png, Inches(0.7), Inches(1.5), Inches(8.6), Inches(5.2))
    _bullets(s, [
        ("Natural partner wins", 0, dict(bold=True, color=INK)),
        ("for 3/4 targets: MDM2, Bcl-xL, 14-3-3.", 0, {}),
        ("CypA is the exception", 0, dict(bold=True, color=INK)),
        ("design -39.6 vs nat -30.8.", 0, {}),
        ("(but see slide 6).", 0, dict(italic=True, color=WARN)),
    ], Inches(9.5), Inches(1.9), Inches(3.5), Inches(4.0), size=14, gap=4)
    _footnote(s, "Error bars are honest +-5 kcal/mol, not the MMPBSA-printed 0.5 to 0.9 "
                 "(frames autocorrelated, n=1). Differences within ~5 kcal/mol are noise. "
                 "Bcl-xL natural ref not converged at 10 ns.")

    # ---- Slide 5: Specificity heatmap (centrepiece) ----
    s = _blank(prs)
    _title(s, "Negative-control specificity matrix",
           "The centrepiece: does the score discriminate sequence and target identity?")
    _pic_fit(s, heat_png, Inches(0.5), Inches(1.5), Inches(7.6), Inches(5.2))
    _bullets(s, [
        ("3/4 targets discriminate", 0, dict(bold=True, color=GOOD)),
        ("both sequence (Ala-scan, scrambled collapse +23 to +44) and target (wrong-target ddG +23 to +64; Bcl-xL and 14-3-3 reach near-zero absolute dG, MDM2 anecdotal due to high SE).", 0, {}),
        ("Pipeline is not just rewarding any peptide in a pocket: sidechains and target identity both matter.", 0, dict(italic=True)),
        ("CypA scrambled does NOT collapse (+1.2).", 0, dict(bold=True, color=WARN)),
        ("That single orange cell is the honesty finding (next slide).", 0, dict(italic=True, color=WARN)),
    ], Inches(8.3), Inches(1.7), Inches(4.7), Inches(4.8), size=13, gap=5)
    _footnote(s, "All cells +-5 kcal/mol. MDM2 wrong-target had high SE (1.3, did not equilibrate in foreign pocket): "
                 "anecdotal. MDM2 scrambled overstates effect (CYX->ALA bug dropped a disulfide).")

    # ---- Slide 6: Honesty finding (CypA) ----
    s = _blank(prs)
    _title(s, "The honesty finding: CypA is composition-driven",
           "A surprising negative result, reported as found.")
    _card(s, Inches(0.7), Inches(1.7), Inches(11.9), Inches(1.5), fill=RGBColor(0xF7, 0xEC, 0xE3))
    tb, tf = _textbox(s, Inches(0.95), Inches(1.85), Inches(11.4), Inches(1.2))
    _set(tf.paragraphs[0], "CypA: a randomly permuted version of the design scores identically to the design "
         "(scrambled ddG = +1.2 kcal/mol, within noise).", 17, WARN, bold=True)
    p = tf.add_paragraph()
    _set(p, "The CypA win is composition-driven shape complementarity, not a sequence-specific binding hypothesis. "
         "The cognate dG is real, but the design's specific sequence contributes essentially nothing.", 14, SLATE)
    _bullets(s, [
        ("Bcl-xL is partially composition-driven: scrambled retains ~60% of binding (scrambled +24.8 vs Ala-scan +37.1).", 0, dict(color=INK)),
        ("MDM2 and 14-3-3 do show sequence specificity (scrambled collapses hard, +24 and +44).", 0, {}),
        ("Implication: 'pipeline ranks design > scrambled' is target-dependent, not a universal property.", 0, dict(bold=True, color=INK)),
        ("Follow-up queued (most important): repeat CypA scrambled with 3 RNG seeds to test if +1.2 is reproducible.", 0, dict(italic=True)),
    ], Inches(0.7), Inches(3.5), Inches(12.0), Inches(3.0), size=15, gap=8)
    _footnote(s, "n=1 trajectory per cell. The +1.2 could be a one-off; the 3-seed repeat is the decisive test and is not yet run.")

    # ---- Slide 7: cheap score vs physics ----
    s = _blank(prs)
    _title(s, "Does the cheap score predict the physics?",
           "ipTM (Boltz confidence) vs MMGBSA dG. The question we cannot yet answer.")
    _bullets(s, [
        ("We would like ipTM to rank designs the way MMGBSA does, so the cheap stage can pre-screen.", 0, dict(color=INK)),
        ("Right now we cannot say whether it does.", 0, dict(bold=True, color=WARN)),
        ("Only 3 designs per target went to MD (12 total in the overnight batch; no alternates). Too few for a per-target Spearman ipTM-vs-MMGBSA.", 0, {}),
        ("ipTM range across the MDM2 designs is narrow (0.83 to 0.96): little spread to correlate against.", 0, {}),
        ("With n this small, any Spearman rho is dominated by noise. Reporting one would be over-claiming.", 0, dict(italic=True)),
        ("Honest verdict: under-powered. The correlation question is deferred to the stratified experiment (slide 8).", 0, dict(bold=True, color=INK)),
    ], Inches(0.7), Inches(1.8), Inches(12.0), Inches(4.5), size=16, gap=10)
    _footnote(s, "No ipTM-MMGBSA correlation is reported because n per target is far too small (3 to 5) and the ipTM "
                 "spread is narrow. An honest rho needs the n=15 stratified design (queued).")

    # ---- Slide 8: What is unclear / queued ----
    s = _blank(prs)
    _title(s, "What is still unclear, and what is queued",
           "Designed and submitted on Slurm. Results pending, no outcomes yet.")
    # two cards: queued-and-running vs needed-not-yet
    _card(s, Inches(0.6), Inches(1.7), Inches(6.0), Inches(4.6), fill=RGBColor(0xE8, 0xF0, 0xE9))
    tb, tf = _textbox(s, Inches(0.85), Inches(1.85), Inches(5.5), Inches(4.3))
    _set(tf.paragraphs[0], "Queued on Slurm (running, results pending)", 16, GOOD, bold=True)
    for txt in [
        "cPEPmatch fit-RMSD sensitivity (MDM2/p53): does seed geometric quality propagate to final design quality, or does ProteinHunter rescue anything? Decision criterion fixed in advance.",
        "ProteinHunter convergence-at-100 (match18_3av9): does the LigandMPNN <-> Boltz loop reach a self-consistent fixed point, or is it a greedy ipTM hill-climb (the 2024 anti-pattern)? 4 independent seeds.",
    ]:
        p = tf.add_paragraph(); p.space_after = Pt(8)
        _set(p, "- " + txt, 12.5, SLATE)
    p = tf.add_paragraph(); p.space_before = Pt(6)
    _set(p, "No results yet. Outcomes not fabricated here.", 11.5, MUTED, italic=True)

    _card(s, Inches(6.9), Inches(1.7), Inches(6.0), Inches(4.6), fill=RGBColor(0xF1, 0xF2, 0xF6))
    tb, tf = _textbox(s, Inches(7.15), Inches(1.85), Inches(5.5), Inches(4.3))
    _set(tf.paragraphs[0], "Needed for an honest verdict (not yet run)", 16, ACCENT, bold=True)
    for txt in [
        "n=15 designs per target with stratified ipTM sampling: the only honest way to estimate Spearman rho for ipTM vs MMGBSA.",
        "100 ns production runs for the per-target winners + natural refs (bronto cannot deliver these; needs the firefly cluster).",
        "Block-averaged / proper standard errors to replace the misleading MMPBSA sqrt-of-frames SE.",
        "MMPBSA-proper (PB, not GB-only) as a second independent solvation model on the winners.",
    ]:
        p = tf.add_paragraph(); p.space_after = Pt(7)
        _set(p, "- " + txt, 12.5, SLATE)
    _footnote(s, "Slurm jobs are designed and submitted; no numbers exist yet. Everything in the right card is identified but unrun.")

    # ---- Slide 9: Limitations & epistemics ----
    s = _blank(prs)
    _title(s, "Limitations and epistemics",
           "Where the numbers are soft, stated plainly.")
    _bullets(s, [
        ("n = 1 trajectory per MD cell. The printed MMPBSA SE (0.5 to 0.9) assumes independent frames; true SE is roughly 3 to 10x larger.", 0, dict(color=INK)),
        ("Treat every per-cell error bar as +-3 to 5 kcal/mol. Any ddG within ~5 of zero is noise.", 0, dict(bold=True, color=WARN)),
        ("Force-field artefacts: 3odi cyclosporin reference unusable (6/11 N-methylated residues defeat GAFF); excluded, not a result.", 0, {}),
        ("Bcl-xL natural reference (Bad) not converged at 10 ns (RMSD still rising, dG drifts ~4.8 kcal/mol): an under-confident floor, recompute at 100 ns.", 0, {}),
        ("MDM2 scrambled control degraded by a CYX->ALA mutagenesis bug (lost a disulfide), so its ddG overstates the sequence effect.", 0, {}),
        ("Scoring is GB-only (igb=5), not PB. Wrong-target pose is a heuristic pocket pushout, not a re-dock (orientation arbitrary).", 0, {}),
        ("Upstream design (per 2026-04-23 critique): Stage 3 optimizes Boltz ipTM, a confidence proxy, not a thermodynamic quantity.", 0, dict(italic=True, color=MUTED)),
    ], Inches(0.7), Inches(1.7), Inches(12.0), Inches(5.0), size=14, gap=8)
    _footnote(s, "These are not hedges added for the talk; each is sourced from the negative-control final report Caveats section.")

    # ---- Slide 10: Conclusion ----
    s = _blank(prs)
    _title(s, "Verdict and recommendation",
           "Honest, not promotional.")
    _card(s, Inches(0.7), Inches(1.65), Inches(11.9), Inches(1.65), fill=RGBColor(0xEC, 0xF1, 0xF5))
    tb, tf = _textbox(s, Inches(0.95), Inches(1.8), Inches(11.4), Inches(1.4))
    _set(tf.paragraphs[0],
         "The pipeline is internally consistent and discriminating, but the designs are NOT yet competitive "
         "with natural binders on dG.", 18, INK, bold=True)
    p = tf.add_paragraph()
    _set(p, "Natural partners beat the designs for 3/4 targets; the one apparent CypA win is composition-driven, "
         "not sequence-specific.", 14, SLATE)
    _bullets(s, [
        ("What is solid: the negative-control matrix shows the score discriminates sequence and target for MDM2, Bcl-xL, 14-3-3.", 0, dict(color=GOOD, bold=True)),
        ("What is not: absolute affinity vs natural binders, the CypA specificity failure, and whether ipTM predicts MMGBSA (under-powered).", 0, dict(color=WARN, bold=True)),
        ("Recommendation: before scaling up, (1) resolve the cheap-score-vs-physics question with the n=15 stratified runs, "
         "(2) get 100 ns + proper SE on the winners via firefly, (3) treat CypA as a cautionary case, not a hit.", 0, dict(color=INK)),
        ("Strategic context (critique 2026-04-23): the interface-mimicry premise is defensible and testable; "
         "the current operationalization (greedy ipTM, GB-only screen) is the weak link, not the idea.", 0, dict(italic=True)),
    ], Inches(0.7), Inches(3.6), Inches(12.0), Inches(3.2), size=14, gap=9)
    _footnote(s, "Verdict rests on n=1 MD with +-5 kcal/mol bars and one un-converged reference. "
                 "Direction is robust; exact magnitudes are not.")


def main():
    data = load_data()
    (ctrl, cog, winners, design_dg, nat_dg, nat_label, nat_flag, ddg) = data
    bar_png = make_bar_chart(design_dg, nat_dg, nat_label, nat_flag)
    heat_png = make_heatmap(ddg)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    build(prs, data, bar_png, heat_png)
    prs.save(OUT_PPTX)

    size_kb = OUT_PPTX.stat().st_size / 1024
    n_slides = len(prs.slides._sldIdLst)
    print(f"Wrote {OUT_PPTX}  ({size_kb:.0f} KB, {n_slides} slides)")
    print(f"Assets: {bar_png.name}, {heat_png.name}")


if __name__ == "__main__":
    main()
